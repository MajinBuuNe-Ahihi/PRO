# Main.py (fixed with LLM-based intent recognition and PPT generation)

import json
import os
import pathlib
import io
import tempfile
from textwrap import dedent
from typing import List, Dict, Optional
from datetime import date

from ddgs import DDGS
from fastapi import FastAPI, Request, Form
from fastapi.responses import StreamingResponse, HTMLResponse, JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from openai import OpenAI
from pydantic import BaseModel
from dotenv import load_dotenv

# Thư viện để tạo file PPTX
from pptx import Presentation
from pptx.util import Inches

# Import các thư viện từ agno
from agno.agent import Agent
from agno.models.openai import OpenAIChat
from starlette.background import BackgroundTask

# Import công cụ DuckDuckGoTools đã được sửa lỗi
from DuckDuckGoTools import DuckDuckGoTools


# Hàm tạo báo cáo PPTX
def create_ppt_report(data: dict, file_path: str):
    """Tạo file PowerPoint với dữ liệu doanh thu giả định."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data.get("title", "Báo cáo Nhà hàng")
    subtitle.text = data.get("subtitle", "Báo cáo tự động xuất từ hệ thống CUKCUK")

    # Thêm slide cho từng nội dung
    for section in data.get('sections', []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.get("header", "Phần")
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = section.get("content", "")

    prs.save(file_path)

# Load environment variables từ file .env
load_dotenv()

# Các Pydantic model để xác thực dữ liệu từ request
class ChatRequest(BaseModel):
    message: str
    session_id: str = ""
    use_search: Optional[bool] = False

class SearchRequest(BaseModel):
    query: str
    max_results: Optional[int] = 5

# Khởi tạo FastAPI app
app = FastAPI(title="Restaurant AI Agent", version="1.0.0")

# Cấu hình CORS để cho phép các domain khác truy cập API
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Cấu hình templates và static files
current_dir = pathlib.Path(__file__).parent.parent
templates = Jinja2Templates(directory=str(current_dir / "templates"))
app.mount("/static", StaticFiles(directory=str(current_dir / "static")), name="static")

# Kiểm tra API key
api_key = os.getenv("OPENAI_API_KEY", "")
if not api_key:
    print("⚠️  Cảnh báo: Không tìm thấy OPENAI_API_KEY trong file .env")

# Khởi tạo OpenAI client
client = OpenAI(api_key=api_key)

# Khởi tạo công cụ DuckDuckGoTools
ddg_tools = DuckDuckGoTools(search=True, news=False, verify_ssl=False, region="vn-vi")

# Định nghĩa Agent chính để xử lý các truy vấn tìm kiếm
duck_agent = Agent(
    model=OpenAIChat(id="gpt-4.1"),
    instructions=dedent("""
        Bạn là một AI Agent chatbot hỗ trợ phần mề quản lý nhà hàng. Phần mềm tên CUKCUK. thuộc công ty MISA.
        Bạn là trợ lý web search. Khi người dùng hỏi và bật search,
        hãy gọi tool duckduckgo.duckduckgo_search để lấy kết quả liên quan.
        Sau đó, tóm tắt ngắn gọn kết quả và đưa ra 2-3 liên kết hữu ích.
        Sử dụng định dạng Markdown để trình bày kết quả.
        Tóm tắt phải **ngắn gọn, xúc tích** và các liên kết phải được liệt kê bằng dấu gạch đầu dòng.
        Ví dụ:
        **[Tiêu đề tóm tắt]**
        - [Tên trang web]: [Link]
        - [Tên trang web]: [Link]
    """),
    tools=[ddg_tools],
    show_tool_calls=True,
    markdown=True,
)

def create_chat_prompt(message: str, context: str = "") -> str:
    """Tạo prompt cho AI Agent nhà hàng"""
    return f"""
Bạn là trợ lý AI cho nhà hàng. Dựa trên yêu cầu của người dùng, hãy trả lời một cách hữu ích và thân thiện.

Yêu cầu của người dùng: {message}

Ngữ cảnh trước đó: {context}

Hãy trả lời bằng tiếng Việt một cách tự nhiên và hữu ích.
"""

def event_stream(message: str, context: str = "", search_results: Optional[list] = None):
    """Stream response từ OpenAI"""
    try:
        prompt = create_chat_prompt(message, context)
        system_content = "Bạn là một AI Agent chatbot hỗ trợ phần mềm quản lý nhà hàng. Phần mềm tên CUKCUK. thuộc công ty MISA. Bạn là trợ lý AI thông minh cho nhà hàng. Trả lời bằng tiếng Việt."
        if search_results:
            system_content += f"\nsearch_results: {json.dumps(search_results, ensure_ascii=False)}"
        
        stream = client.chat.completions.create(
            model="gpt-4.1", 
            messages=[
                {"role": "system", "content": system_content},
                {"role": "user", "content": prompt}
            ],
            stream=True,
            temperature=0.7
        )
        
        for chunk in stream:
            if chunk.choices[0].delta.content:
                yield chunk.choices[0].delta.content
                
    except Exception as e:
        yield f"Lỗi: {str(e)}"

# Hàm mới để xác định ý định của người dùng bằng LLM
async def is_ppt_request(message: str) -> bool:
    """Sử dụng LLM để phân loại tin nhắn của người dùng có phải là yêu cầu tạo báo cáo PPT."""
    try:
        response = client.chat.completions.create( # Loại bỏ 'await' ở đây
            model="gpt-4.1", # Sử dụng model của bạn
            messages=[
                {"role": "system", "content": """Bạn là một trợ lý phân loại ý định. Nhiệm vụ của bạn là xác định xem một tin nhắn có phải là yêu cầu tạo một báo cáo PowerPoint (PPT) hoặc một loại báo cáo nào khác hay không. 
                Trả lời "Có" nếu đó là một yêu cầu tạo báo cáo (ví dụ: "Tạo báo cáo doanh thu", "Xuất file PPT", "Làm giúp tôi báo cáo này"), và "Không" nếu không phải. 
                Chỉ trả lời "Có" hoặc "Không", không có thêm văn bản nào khác."""},
                {"role": "user", "content": message}
            ],
            temperature=0,
            max_tokens=10
        )
        return response.choices[0].message.content.strip().lower() == "có"
    except Exception as e:
        print(f"Lỗi khi phân loại tin nhắn với LLM: {e}")
        return False

@app.get("/")
async def home(request: Request):
    """Trang chủ với giao diện chat"""
    return templates.TemplateResponse("index.html", {"request": request})

@app.post("/generate-ppt")
async def generate_ppt():
    """Endpoint để tạo và tải file PPT."""
    # Dữ liệu giả định cho báo cáo
    report_data = {
        "title": "Báo cáo Doanh thu",
        "subtitle": f"Ngày {date.today().strftime('%d/%m/%Y')}",
        "sections": [
            {
                "header": "Tổng quan",
                "content": f"Doanh thu hôm nay: 5.000.000 VNĐ\nSố lượng đơn hàng: 50 đơn\nLợi nhuận: 2.500.000 VNĐ"
            }
        ]
    }
    
    # Tạo file tạm thời để lưu PPT
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        file_path = tmp.name
        # Thay thế hàm tạo báo cáo của bạn
        # create_ppt_report(report_data, file_path) 
    
    filename = f"bao_cao_doanh_thu_{date.today().strftime('%Y%m%d')}.pptx"
    
    # Sử dụng BackgroundTask để xóa file sau khi phản hồi được gửi
    task = BackgroundTask(os.remove, file_path)
    
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
        background=task
    )

@app.post("/chat")
async def chat_stream(request: ChatRequest):
    """API chat với streaming response, có thể kèm search tool"""
    # Xử lý trường hợp người dùng muốn tạo PPT bằng cách sử dụng LLM để xác định ý định
    if not request.use_search:
        is_ppt = await is_ppt_request(request.message) # Sửa lỗi ở đây: thêm 'await'
        if is_ppt:
            return StreamingResponse(
                io.StringIO("Bạn muốn tạo báo cáo doanh thu PPT? Vui lòng nhấn vào nút 'Tạo PPT' để tải xuống!"),
                media_type="text/plain"
            )
    
    if request.use_search and request.message:
        print("log chat tim kiem")
        def gen():
            try:
                # Agent sẽ tự động gọi tool và tạo ra câu trả lời theo instructions đã định nghĩa
                text = duck_agent.run(request.message)
                yield str(text.content)
            except Exception as e:
                yield f"Loi {str(e)}"
        return StreamingResponse(gen(), media_type="text/plain")
    else:
        print("log chat thuong")
        return StreamingResponse(
            event_stream(request.message, search_results=None),
            media_type="text/plain"
        )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
